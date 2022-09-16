import os
from re import L
import requests
import json
import uuid
from io import BytesIO

import streamlit as st
from streamlit.web.server import Server

from pptx import Presentation
import extra_streamlit_components as stx

st.set_page_config(page_title="tools", page_icon="https://www.3drepo.io/favicon.ico", layout="centered", initial_sidebar_state="auto", menu_items=None)
st.header("3D Repo Safetibase PowerPoint App")

@st.cache(allow_output_mutation=True)
def get_manager():
    return stx.CookieManager()

def selfGet(url):
    curSession = requests.Session() 
    if connectsid:
        cookie = {
            'connect.sid' : connectsid,
        }
        return curSession.get(url, cookies=cookie)
    else:
        url = url + "?" + needKey()
        return curSession.get(url)

def get_3dreporisks(domain,teamspace,model,api_key):
    url = domain + "/api/"+teamspace+"/"+model+"/risks"
    risk_response = selfGet(url)
    risk_response_object = json.loads(risk_response.text)
    return risk_response_object

def get_3drepologin(domain, connectsid):
    url = domain + "/api/me"
    login_response = selfGet(url)
    return login_response

def insert(domain,teamspace,model,apiKey,output):

    if not teamspace and model:
        return False
        
    SLD_LAYOUT_TITLE_AND_CONTENT = 0

    prs = Presentation('Template.pptx')
    slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]

    risks = get_3dreporisks(domain,teamspace,model,apiKey)

    if 'status' in risks:
        if risks['status'] == 401:
            st.text(risks)
            return False
    
    if risks:
        for risk in risks:
            slide = prs.slides.add_slide(slide_layout)
            name = slide.placeholders[0]
            name.text = risk['name']

            url = slide.placeholders[15]
            urlText = domain + "/viewer/" + teamspace + "/" + model + "?risk=" + risk['_id'] 
            url.text = urlText
            try:
                desc = slide.placeholders[13]
                desc.text = risk['desc']
                image = slide.placeholders[14]
            except:
                continue
            try:
                imageLocation = domain + "/api/" + risk['viewpoint']['screenshot'] + "?" + needKey()
                imageGet = requests.get(imageLocation)
                filename = str(uuid.uuid4())
                file = open(filename, "wb")
                file.write(imageGet.content)
                file.close()
                image.insert_picture(filename)
            except:
                continue
            try:
                os.remove(filename)
            except:
                continue
        fileName = output + '.pptx'
        outputFile = BytesIO()
        prs.save(outputFile)
        outputFile.seek(0)

        st.download_button(
            label="Download data as Powerpoint",
            data=outputFile,
            file_name=fileName
        )
    else:
        st.write("No risks found chosen container.")

connectsid = ''

if "DEPLOY_API" in os.environ:
    domainValue = "https://" + os.environ["DEPLOY_API"]
else:
    domainValue = "https://www.3drepo.io"

domain = st.text_input("Domain:", value = domainValue)
output = st.text_input("Output File Name:", value = "3D Repo Safetibase Export")

def loadProjects(teamspace):
    url = domain + "/api/" + teamspace + "/projects"
    projects = selfGet(url)
    try:
        return json.loads(projects.text)
    except Exception as err:
        st.write(projects.text)
        st.write(projects)
        st.write(err)

def loadModels(teamspace,project):
    url = domain + "/api/" + teamspace + "/projects/" + project + "/models"
    models = selfGet(url)
    try:
        return json.loads(models.text)
    except Exception as err:
        st.write(models)
        st.write(err)

def loadAll(teamspace):
    url = domain + "/api/" + teamspace + ".json"
    all = selfGet(url)
    try:
        return json.loads(all.text)
    except Exception as err:
        st.write(all)
        st.write(err)


def needKey():
    if apiKey:
        return "key=" + apiKey
    else:
        return ""

if not st.experimental_user.email == 'test@localhost.com':
    connectsid = st.experimental_user.email
    login_response = get_3drepologin(domain,connectsid)
    login_response_success = login_response.status_code == 200
    if login_response_success:
        teamspace = login_response.json()['username']
        everything = loadAll(teamspace)
        teamspaces = {}
        for id,accounts in enumerate(everything['accounts']):
            teamspaces[accounts['account']] = id
    else:
        everything = []
else:
    everything = []
    login_response_success = False

if everything:
    st.session_state.current_teamspace = st.selectbox(
        'Which teamspace are we interested in?',
        teamspaces.keys()
        )
else:
    st.session_state.current_teamspace = ""

if st.session_state.current_teamspace:
    projectsList = {}
    for id,projects in enumerate(everything['accounts'][teamspaces[st.session_state.current_teamspace]]['projects']):
        projectsList[projects['name']] = id

    st.session_state.current_project = st.selectbox(
        'Which project are we interested in?',
        projectsList.keys()
        )
else:
    st.session_state.current_project = ""

if st.session_state.current_project:
    modelsList = {}
    for id,models in enumerate(everything['accounts'][teamspaces[st.session_state.current_teamspace]]['projects'][projectsList[st.session_state.current_project]]['models']):
        modelsList[models['name']] = models['model']
    st.session_state.current_model = st.selectbox(
        'Which container are we interested in?',
        modelsList.keys()
        )
else:
    st.session_state.current_model = ""

if not login_response_success:
    apiKey = st.text_input("API Key:")
else:
    st.text("Logged in as : " + login_response.json()['username'])
    apiKey = ''

if apiKey:
    connectsid = ''


if not st.session_state.current_teamspace:
    st.text("Insert your Model/Federation details below to generate a Powerpoint file of all the Risks")

if st.session_state.current_teamspace:
    teamspace = st.session_state.current_teamspace
else:
    teamspace = st.text_input("Teamspace:")

if st.session_state.current_model:
    model = modelsList[st.session_state.current_model]
else:
    model = st.text_input("Model:")

if st.button("Submit"):
    insert(domain,teamspace,model,apiKey,output)

if "DEPLOY_TAG" in os.environ:
    st.write(os.environ["DEPLOY_TAG"])
